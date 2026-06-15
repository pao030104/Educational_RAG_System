"""
PowerPoint 文档加载器模块
===========================
基于 python-pptx 的 PPT/PPTX 加载器，提取幻灯片中的文本和图片文字。

工作流程:
    ┌──────────────┐    ┌──────────────────┐    ┌──────────────────┐
    │ 打开 PPT 文件  │───→│ 逐幻灯片处理      │───→│ 按位置排序形状     │
    │              │    │ (按从上到下从左   │    │ (top, left 坐标)  │
    │              │    │  到右顺序阅读)    │    │                  │
    └──────────────┘    └──────────────────┘    └──────────────────┘
                                                       │
                                                       ▼
    ┌──────────────┐    ┌──────────────────┐    ┌──────────────────┐
    │ 合并输出      │←───│ OCR 识别图片文字   │←───│ 分类处理:        │
    │              │    │                  │    │ • 文本框 → 直接取 │
    └──────────────┘    └──────────────────┘    │ • 表格 → 遍历提取 │
                                                │ • 图片 → OCR 识别 │
                                                │ • 组合 → 递归处理 │
                                                └──────────────────┘

形状类型码:
    13: 图片 (Picture)
    6:  组合 (Group) - 需要递归处理子形状
    其他: 文本框、表格、图表等（通过 has_text_frame / has_table 判断）
"""

from typing import Iterator                        # 类型注解
from edu_ocr import get_ocr                        # OCR 引擎
from langchain_core.documents import Document       # LangChain Document
from langchain_core.document_loaders import BaseLoader  # 基类
from pptx import Presentation                       # PPT 文件解析
from PIL import Image                               # 图片处理
import numpy as np                                  # 数组转换
from io import BytesIO                              # 字节流转图片
from tqdm import tqdm                                # 进度条


class OCRPPTLoader(BaseLoader):
    """
    PowerPoint 文档加载器，提取文本框、表格和图片中的全部文字。

    遵循 LangChain BaseLoader 接口。

    属性:
        filepath (str): PPT/PPTX 文件的路径。

    使用示例:
        >>> loader = OCRPPTLoader(filepath="/path/to/slides.pptx")
        >>> docs = loader.load()
        >>> for doc in docs:
        ...     print(doc.page_content)
    """

    def __init__(self, filepath: str) -> None:
        """
        初始化 PPT 加载器。

        参数:
            filepath (str): PPT/PPTX 文件的绝对或相对路径。
        """
        self.filepath = filepath

    def lazy_load(self) -> Iterator[Document]:
        """
        延迟加载模式：提取文本后返回单个 Document。

        Returns:
            Iterator[Document]
        """
        line = self.ppt2text(self.filepath)
        yield Document(page_content=line, metadata={"source": self.filepath})

    def ppt2text(self, filepath):
        """
        从 PPT 文件提取全部文本内容。

        处理内容类型:
            - 文本框 (Text Frame): 直接读取 shape.text
            - 表格 (Table): 遍历行→单元格→段落
            - 图片 (Picture, shape_type=13): OCR 识别
            - 组合形状 (Group, shape_type=6): 递归处理子形状

        阅读顺序:
            幻灯片内的形状按 top (从上到下) 和 left (从左到右) 排序，
            模拟人类的自然阅读顺序。

        参数:
            filepath (str): PPT 文件路径。

        返回:
            str: 提取的全部文本，各部分以换行分隔。
        """
        # 打开 PowerPoint 文件
        prs = Presentation(filepath)
        print(f'prs-->{prs}')

        # 获取 OCR 引擎
        ocr = get_ocr()
        resp = ""  # 累积全部文本

        def extract_text(shape):
            """
            递归提取形状中的文本（内部函数）。

            使用 nonlocal 声明访问和修改外部函数中的 resp 变量，
            避免每次递归都传递 resp。

            参数:
                shape: python-pptx 的形状对象。
            """
            nonlocal resp  # 声明使用外部函数的 resp 变量

            # ---- 处理文本框 ----
            # has_text_frame: 检查形状是否包含文本框/占位符
            if shape.has_text_frame:
                resp += shape.text.strip() + "\n"

            # ---- 处理表格 ----
            # has_table: 检查形状是否为表格
            if shape.has_table:
                for row in shape.table.rows:           # 遍历所有行
                    for cell in row.cells:              # 遍历每行的单元格
                        for paragraph in cell.text_frame.paragraphs:  # 遍历段落
                            resp += paragraph.text.strip() + "\n"

            # ---- 处理图片 ----
            # shape_type == 13 表示这是一个 Picture 对象
            if shape.shape_type == 13:
                # shape.image.blob: 图片的原始二进制数据
                image = Image.open(BytesIO(shape.image.blob))
                # OCR 识别（将 PIL Image 转为 numpy 数组）
                result, _ = ocr(np.array(image))
                if result:
                    ocr_result = [line[1] for line in result]
                    resp += "\n".join(ocr_result)

            # ---- 处理组合形状（递归） ----
            # shape_type == 6 表示这是一个 GroupShape
            elif shape.shape_type == 6:
                # 递归处理组合中的每个子形状
                for child_shape in shape.shapes:
                    extract_text(child_shape)

        # ---- 创建进度条 ----
        b_unit = tqdm(
            total=len(prs.slides),
            desc="OCRPPTLoader slide index: 1"
        )

        # ---- 逐幻灯片处理 ----
        for slide_number, slide in enumerate(prs.slides, start=1):
            # 更新进度条
            b_unit.set_description(
                "OCRPPTLoader slide index: {}".format(slide_number)
            )
            b_unit.refresh()

            # 按位置排序: 先按 top (从上到下)，再按 left (从左到右)
            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (x.top, x.left)
            )

            # 提取每个形状的文本
            for shape in sorted_shapes:
                extract_text(shape)

            b_unit.update(1)

        return resp


if __name__ == '__main__':
    # 测试 PPT 加载器
    # 请将路径修改为你本地的 PPT 文件路径
    import sys
    test_file = sys.argv[1] if len(sys.argv) > 1 else './samples/ocr_01.pptx'
    img_loader = OCRPPTLoader(filepath=test_file)
    doc = img_loader.load()
    print(doc)
